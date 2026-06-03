"""POST /api/science/chat — SSE streaming endpoint for external backends.
用于科研助手智能体
Accepts the standardised ModelChat request format, creates a Hermes
chat completion, subscribes to the SSE event stream, and re-emits events
in the format expected by the downstream consumer (Yang's backend / front-end).

Uses /v1/chat/completions with X-Hermes-Session-Id header so the Hermes
container automatically loads conversation history from its SQLite session DB.

Response SSE payload types:
  4 = text stream delta (message for final text, reasoningMessage for thinking)
"""

from __future__ import annotations

import asyncio
import io
import json
import logging
import mimetypes
import os
import re
import tarfile
from typing import Any
from urllib.parse import quote, urlparse

import httpx
import requests as _sync_requests
from docker.errors import APIError as DockerAPIError
from docker.errors import NotFound as DockerNotFound
from fastapi import APIRouter, HTTPException, Query, Request
from fastapi.responses import Response, StreamingResponse
from pydantic import BaseModel, field_validator

from app.config import settings
from app.container.manager import get_docker_container
from app.training_trace import append_jsonl_trace, build_model_chat_trace_record

logger = logging.getLogger(__name__)
router = APIRouter(tags=["model-chat"])

# ---------------------------------------------------------------------------
# Request / Response models
# ---------------------------------------------------------------------------


class ModelChatMessage(BaseModel):
    role: str  # system | user | assistant | tool
    content: Any = ""  # str for plain text, or list[dict] for multi-content (text + file attachments)


class ModelChatRequest(BaseModel):
    linkId: str
    sessionId: str
    userId: int = 1
    functionId: int = 1
    messages: list[ModelChatMessage] = []
    type: int = 0  # -1 = abort
    attachment: Any = {}
    callTools: bool = True
    XAPIVersion: Any = 1

    @field_validator("attachment", mode="before")
    @classmethod
    def _coerce_attachment(cls, v: Any) -> dict[str, Any]:
        if isinstance(v, dict):
            return v
        return {}

    @field_validator("XAPIVersion", mode="before")
    @classmethod
    def _coerce_xapi_version(cls, v: Any) -> int:
        try:
            return int(v) if v is not None else 1
        except (TypeError, ValueError):
            return 1


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

# Session → run_id mapping so we can abort by sessionId
_session_run_map: dict[str, str] = {}
_session_run_cleanup_tasks: dict[str, asyncio.TimerHandle] = {}
# Session → container URL (for session affinity, file download, abort routing)
_session_container_map: dict[str, str] = {}

# Container URL → session_id (tracks which session currently holds the container)
_busy_containers: dict[str, str] = {}

# Async lock protecting _busy_containers, _rr_index, and writes to _session_container_map
_container_lock = asyncio.Lock()

# Explicit round-robin index as mutable container (replaces itertools.cycle)
_rr_index = [0]

_TYPEWRITER_DELAY = 0.05  # seconds between tokens
_TYPEWRITER_FAST_DELAY = 0.01  # when queue > 20
_SESSION_CONTAINER_TTL = 30 * 60  # 30 minutes before cleaning up session→container mapping
_STREAM_REWRITE_TAIL_CHARS = 512  # keep enough tail text for split file paths


def _build_base(req: ModelChatRequest) -> dict[str, Any]:
    return {
        "linkId": req.linkId,
        "sessionId": req.sessionId,
        "userId": req.userId,
        "functionId": req.functionId,
        "attachment": req.attachment or {},
        "XAPIVersion": req.XAPIVersion,
    }


def _sse_line(data: dict[str, Any]) -> str:
    return f"data: {json.dumps(data, ensure_ascii=False)}\n\n"


def _format_tool_event(event: dict[str, Any]) -> str | None:
    """Convert a hermes tool event to a human-readable progress string."""
    evt = event.get("event") or event.get("type") or ""
    if evt == "tool.started":
        name = event.get("tool") or "unknown"
        return f"[工具] 调用: {name}"
    if evt == "tool.completed":
        name = event.get("tool") or "unknown"
        err = " (失败)" if event.get("error") else ""
        return f"[工具] 完成: {name}{err}"
    return None


def _container_name_from_url(url: str) -> str:
    """Extract container name from URL like 'http://hermes-innovation-03:18080'."""
    from urllib.parse import urlparse
    return urlparse(url).hostname or ""


# Regex to match file paths in agent messages (absolute paths with common extensions)
_FILE_EXTENSIONS = r'(?:md|docx|pdf|pptx|xlsx|csv|html|txt|json|png|jpg|jpeg|gif|svg|zip|tar|gz)'
# Pattern 1: backtick-wrapped paths like `/tmp/report.pdf`
_FILE_PATH_BACKTICK_RE = re.compile(
    r'`(/(?!/)[^`\s]+\.' + _FILE_EXTENSIONS + r')`',
)
# Pattern 2: bare paths (preceded by whitespace, colon, or line start; not URLs)
_FILE_PATH_BARE_RE = re.compile(
    r'(?:^|[\s：:*])(/(?!/)[^\s\'"<>]+\.' + _FILE_EXTENSIONS + r')',
    re.MULTILINE,
)


def _inject_download_urls(text: str, session_id: str) -> str:
    """Replace file paths in text with markdown download links.

    Two passes: backtick-wrapped paths first, then bare paths.
    """
    def _make_link(path: str) -> str:
        filename = path.rsplit("/", 1)[-1]
        base = os.environ.get("DOWNLOAD_BASE_URL", "").rstrip("/")
        url = f"{base}/api/science/chat/file?sessionId={quote(session_id)}&path={quote(path)}"
        return f'[{filename}]({url})'

    # Pass 1: backtick-wrapped  `/path/file.pdf` → [file.pdf](url)
    def _replace_backtick(m: re.Match) -> str:
        return _make_link(m.group(1))
    text = _FILE_PATH_BACKTICK_RE.sub(_replace_backtick, text)

    # Pass 2: bare paths  /path/file.pdf → [file.pdf](url)
    def _replace_bare(m: re.Match) -> str:
        path = m.group(1)
        prefix = m.group(0)[0] if m.group(0)[0] != "/" else ""
        return f'{prefix}{_make_link(path)}'
    text = _FILE_PATH_BARE_RE.sub(_replace_bare, text)

    return text


# Tokenise text for typewriter: Chinese chars individually, English words together
_TOKEN_RE = re.compile(r"[a-zA-Z]+|\S|\s", re.UNICODE)


def _tokenize(text: str) -> list[str]:
    return _TOKEN_RE.findall(text) or [text]


def _download_and_extract_text(url: str) -> str:
    """Download an attachment file from URL and extract its text content.

    Supports .docx (OOXML zip parsing), .pdf (pypdf), .pptx (python-pptx),
    and plain text files.
    """
    import io as _io
    import zipfile
    from xml.etree import ElementTree

    # --- SSRF 防护：只允许下载白名单域名下的附件 ---
    # 如需新增允许的附件域名（如对象存储/CDN/其他环境），在此元组追加即可。
    _ALLOWED_ATTACHMENT_DOMAINS = (".infox-med.com",)
    try:
        _host = (urlparse(url).hostname or "").lower()
    except Exception:
        _host = ""
    _allowed = any(
        _host == d.lstrip(".") or _host.endswith(d)
        for d in _ALLOWED_ATTACHMENT_DOMAINS
    )
    if not _allowed:
        logger.warning("Blocked attachment download from non-whitelisted host: %s (url=%s)", _host, url)
        return ""

    # 提取文本的字符上限，防止超大附件撑爆上下文/费用
    _MAX_ATTACHMENT_CHARS = 50000

    try:
        resp = _sync_requests.get(url, timeout=30)
        resp.raise_for_status()
    except Exception as exc:
        logger.warning("Failed to download attachment from %s: %s", url, exc)
        return ""

    content_type = resp.headers.get("Content-Type", "")
    content = resp.content

    # Determine likely file type from URL extension
    filename = url.rsplit("/", 1)[-1].lower() if "/" in url else ""
    is_docx = filename.endswith(".docx") or "docx" in content_type

    if is_docx:
        try:
            with zipfile.ZipFile(_io.BytesIO(content)) as z:
                if "word/document.xml" in z.namelist():
                    xml_bytes = z.read("word/document.xml")
                    root = ElementTree.fromstring(xml_bytes)
                    ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                    texts: list[str] = []
                    for t_elem in root.iter(f"{{{ns}}}t"):
                        if t_elem.text:
                            texts.append(t_elem.text)
                    # Join paragraphs with newlines
                    paragraphs: list[str] = []
                    for p_elem in root.iter(f"{{{ns}}}p"):
                        para_texts: list[str] = []
                        for t_elem in p_elem.iter(f"{{{ns}}}t"):
                            if t_elem.text:
                                para_texts.append(t_elem.text)
                        if para_texts:
                            paragraphs.append("".join(para_texts))
                    _docx_text = "\n\n".join(paragraphs) if paragraphs else "\n".join(texts)
                    return _docx_text[:_MAX_ATTACHMENT_CHARS]
        except Exception as exc:
            logger.warning("Failed to extract docx text from %s: %s", url, exc)

    # PDF extraction
    is_pdf = filename.endswith(".pdf") or "pdf" in content_type
    if is_pdf:
        try:
            from pypdf import PdfReader
            reader = PdfReader(_io.BytesIO(content))
            _page_texts: list[str] = []
            for page in reader.pages:
                _t = page.extract_text()
                if _t:
                    _page_texts.append(_t)
            _pdf_text = "\n\n".join(_page_texts)
            if _pdf_text.strip():
                return _pdf_text[:_MAX_ATTACHMENT_CHARS]
            else:
                logger.warning("No extractable text from PDF (possibly scanned/image-based): %s", url)
        except Exception as exc:
            logger.warning("Failed to extract PDF text from %s: %s", url, exc)

    # PPTX extraction
    is_pptx = filename.endswith(".pptx") or "pptx" in content_type
    if is_pptx:
        logger.info("输入的pptx文件")
        try:
            from pptx import Presentation
            prs = Presentation(_io.BytesIO(content))
            _slides: list[str] = []
            for slide in prs.slides:
                _shape_texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            _t = para.text.strip()
                            if _t:
                                _shape_texts.append(_t)
                if _shape_texts:
                    _slides.append("\n".join(_shape_texts))
            _pptx_text = "\n\n---\n\n".join(_slides)
            if _pptx_text.strip():
                logger.info(f"返回的内容是:{_pptx_text}")
                return _pptx_text[:_MAX_ATTACHMENT_CHARS]
            else:
                logger.warning("No extractable text from PPTX: %s", url)
        except Exception as exc:
            logger.warning("Failed to extract PPTX text from %s: %s", url, exc)

    # Image formats — use VL model to recognize content
    _image_extensions = (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff", ".tif")
    is_image = any(filename.endswith(ext) for ext in _image_extensions) or content_type.startswith("image/")
    if is_image:
        logger.info("输入的图片文件")
        try:
            from app.routes.read_image import recognize_image_scene
            success, description = asyncio.run(recognize_image_scene(url, "请描述这张图片的内容"))
            if success and description:
                logger.info("图片识别成功", description)
                return f"[图片内容识别]\n{description}"[:_MAX_ATTACHMENT_CHARS]
        except Exception as exc:
            logger.warning("Failed to recognize image from %s: %s", url, exc)

    # Try plain text
    try:
        text = content.decode("utf-8")
        # Basic check: not purely binary
        if any(0x20 <= b <= 0x7e or b in (0x0a, 0x0d, 0x09) for b in content[:512]):
            return text[:_MAX_ATTACHMENT_CHARS]
    except (UnicodeDecodeError, Exception):
        pass

    logger.warning("Could not extract text from attachment: %s (type=%s)", url, content_type)
    return ""


def _extract_content_text(content: Any) -> str:
    """Extract plain text from a message content field.

    Handles both legacy plain-string format and the multi-content array format
    that includes file attachments.
    """
    logger.info(f"开始提取content, {content}")
    if isinstance(content, str):
        return content

    if not isinstance(content, list):
        return str(content) if content else ""

    parts: list[str] = []
    logger.info(f"提取block开始")
    for block in content:
        logger.info(f"提取block内容:{block}")
        if not isinstance(block, dict):
            continue
        block_type = block.get("type", "")
        if block_type == "text":
            text = block.get("text", "")
            if text:
                parts.append(text)
        elif block_type == "file" or block_type == "image_url":
            url = block.get("url", "")
            if url:
                extracted = _download_and_extract_text(url)
                if extracted:
                    logger.info(f"提取的附件的内容:{extracted}")
                    parts.append(f"\n\n[附件内容]\n{extracted}")
                else:
                    logger.warning("No text extracted from attachment: %s", url)

    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Hermes connection helpers
# ---------------------------------------------------------------------------

def _hermes_auth_headers() -> dict[str, str]:
    key = settings.dedicated_hermes_api_key
    if not key:
        return {}
    return {"Authorization": f"Bearer {key}"}


# Pool of 10 pre-created science hermes containers
_SCIENCE_CONTAINERS = [f"http://hermes-science-{i:02d}:{settings.dedicated_hermes_internal_port}" for i in range(10)]


async def _resolve_hermes_url(session_id: str) -> str:
    """Return the base URL for a pre-created science hermes runtime.

    Session-affinity round-robin: if a session previously used a container and
    that container is not busy, prefer it. Otherwise fall back to round-robin
    across available containers.
    """
    async with _container_lock:
        # 1. Check session affinity: prefer previously assigned container
        preferred = _session_container_map.get(session_id)
        if preferred and preferred not in _busy_containers:
            _busy_containers[preferred] = session_id
            logger.info("[model-chat] 会话亲和: session=%s → %s (复用已分配容器)", session_id, preferred)
            return preferred

        # 2. Scan for a non-busy container starting from round-robin index
        n = len(_SCIENCE_CONTAINERS)
        for offset in range(n):
            idx = (_rr_index[0] + offset) % n
            url = _SCIENCE_CONTAINERS[idx]
            if url not in _busy_containers:
                _busy_containers[url] = session_id
                _session_container_map[session_id] = url
                _rr_index[0] = (idx + 1) % n
                logger.info("[model-chat] 轮询分配: session=%s → %s (preferred=%s, busy=%s)", session_id, url, preferred, bool(preferred))
                return url

        # 3. All busy: fallback to preferred or round-robin pick
        if preferred:
            chosen = preferred
        else:
            chosen = _SCIENCE_CONTAINERS[_rr_index[0]]
            _rr_index[0] = (_rr_index[0] + 1) % n
        _busy_containers[chosen] = session_id
        _session_container_map[session_id] = chosen
        logger.warning("[model-chat] 所有容器均忙碌，复用容器: %s (session=%s)", chosen, session_id)
        return chosen


async def _release_container(container_url: str, session_id: str) -> None:
    """Mark a container as no longer busy after its SSE stream completes."""
    async with _container_lock:
        # Only release if the session still owns this container entry
        if _busy_containers.get(container_url) == session_id:
            _busy_containers.pop(container_url, None)


# ---------------------------------------------------------------------------
# Shared core logic
# ---------------------------------------------------------------------------

async def _abort_session(session_id: str) -> dict[str, Any] | None:
    """Try to abort a running session. Returns the response dict or None."""
    run_id = _session_run_map.get(session_id)
    if not run_id:
        return None
    try:
        base_url = _session_container_map.get(session_id)
        if not base_url:
            base_url = await _resolve_hermes_url(session_id)
            await _release_container(base_url, session_id)
        async with httpx.AsyncClient(timeout=10.0) as client:
            await client.post(
                f"{base_url}/v1/runs/{run_id}/stop",
                headers=_hermes_auth_headers(),
            )
    except Exception as exc:
        logger.warning("abort failed for session %s: %s", session_id, exc)
    _session_run_map.pop(session_id, None)
    return None


async def _do_science_chat(
    request: Request,
    body: ModelChatRequest,
    chat_messages: list[dict[str, str]],
    log_tag: str = "model-chat",
) -> StreamingResponse:
    """Core chat/translate logic shared by all science chat endpoints.

    Parameters
    ----------
    chat_messages : list[dict[str, str]]
        Already-prepared messages to send to the hermes container.
        Callers should have applied any prompt modifications (e.g. translation prefix).
    log_tag : str
        Tag used in log messages to distinguish chat vs translate endpoints.
    """
    # --- Validate ---
    if not body.linkId or not body.sessionId:
        raise HTTPException(status_code=400, detail="linkId and sessionId are required")
    if not body.messages:
        raise HTTPException(status_code=400, detail="messages array is required and must not be empty")

    logger.info(
        "[%s] 请求开始: linkId=%s sessionId=%s userId=%s functionId=%s type=%s",
        log_tag, body.linkId, body.sessionId, body.userId, body.functionId, body.type,
    )

    base_url = await _resolve_hermes_url(body.sessionId)
    base = _build_base(body)
    trace_tool_events: list[dict[str, Any]] = []
    trace_run_id = ""
    trace_status = "not_started"
    trace_written = False

    # Build headers with X-Hermes-Session-Id for automatic history loading
    hermes_headers = _hermes_auth_headers()
    hermes_headers["X-Hermes-Session-Id"] = body.sessionId

    async def _stream():
        nonlocal trace_run_id, trace_status, trace_written
        full_text = ""           # 累积完整消息原文
        pending_text = ""        # 暂存尾部，避免跨 chunk 文件路径被改写坏
        _usage: dict[str, int] = {}   # 从最终 chunk 中提取的 token 用量
        _model_name: str = ""         # 实际使用的模型名

        def _pop_rewritten_pending(force: bool = False) -> str:
            nonlocal pending_text
            if not pending_text:
                return ""
            if force:
                text = pending_text
                pending_text = ""
            elif len(pending_text) > _STREAM_REWRITE_TAIL_CHARS:
                cut = len(pending_text) - _STREAM_REWRITE_TAIL_CHARS
                text = pending_text[:cut]
                pending_text = pending_text[cut:]
            else:
                return ""
            return _inject_download_urls(text, body.sessionId)

        def _trace_enabled() -> bool:
            return bool(getattr(settings, "training_trace_enabled", False))

        def _write_trace_once(status_text: str) -> None:
            nonlocal trace_written
            if trace_written or not _trace_enabled():
                return
            trace_written = True
            try:
                record = build_model_chat_trace_record(
                    link_id=body.linkId,
                    session_id=body.sessionId,
                    request_user_id=body.userId,
                    function_id=body.functionId,
                    messages=body.messages,
                    user=None,
                    run_id=trace_run_id,
                    model="hermes-agent",
                    runtime="hermes",
                    tool_events=trace_tool_events,
                    final_output=full_text,
                    status=status_text,
                    trace_hash_salt=getattr(settings, "training_trace_hash_salt", ""),
                )
                append_jsonl_trace(getattr(settings, "training_trace_dir", ".hermes/training_traces"), record)
            except Exception as exc:
                logger.warning("[%s] training trace write failed: %s", log_tag, exc)

        # 立刻返回 thinking，让用户感受到响应
        yield _sse_line({**base, "message": "", "reasoningMessage": "thinking:", "type": 4})

        # 1. Create chat completion stream
        try:
            client = httpx.AsyncClient(timeout=None)
            req = client.build_request(
                "POST",
                f"{base_url}/v1/chat/completions",
                headers=hermes_headers,
                json={
                    "model": "hermes-agent",
                    "messages": chat_messages,
                    "stream": True,
                },
            )
            event_stream = await client.send(req, stream=True)
        except Exception as exc:
            logger.error("[%s] create chat completion error: %s", log_tag, exc)
            _write_trace_once("create_chat_completion_error")
            await _release_container(base_url, body.sessionId)
            yield _sse_line({**base, "message": "[stop]", "reasoningMessage": "", "type": 4})
            return

        if event_stream.status_code >= 400:
            logger.error("[%s] chat completion returned %s", log_tag, event_stream.status_code)
            _write_trace_once("stream_connect_failed")
            await event_stream.aclose()
            await client.aclose()
            await _release_container(base_url, body.sessionId)
            yield _sse_line({**base, "message": "[stop]", "reasoningMessage": "", "type": 4})
            return

        # 2. Process SSE event stream
        try:
            buffer = ""
            current_event = ""  # for named events like hermes.tool.progress
            async for chunk in event_stream.aiter_bytes():
                if await request.is_disconnected():
                    trace_status = "disconnected"
                    break

                buffer += chunk.decode("utf-8", errors="ignore")
                while "\n\n" in buffer:
                    raw_block, buffer = buffer.split("\n\n", 1)

                    # Parse event type if present (e.g. "event: hermes.tool.progress")
                    for line in raw_block.splitlines():
                        if line.startswith("event:"):
                            current_event = line[6:].strip()

                    # Parse data lines
                    data_lines = [
                        line[5:].strip()
                        for line in raw_block.splitlines()
                        if line.startswith("data:")
                    ]
                    if not data_lines:
                        current_event = ""
                        continue
                    data_str = "\n".join(data_lines)
                    if data_str == "[DONE]":
                        current_event = ""
                        break

                    try:
                        data = json.loads(data_str)
                    except json.JSONDecodeError:
                        current_event = ""
                        continue
                    if not isinstance(data, dict):
                        current_event = ""
                        continue

                    # Capture run_id from first chunk (for abort + training trace)
                    if not trace_run_id:
                        trace_run_id = data.get("id", "") or ""
                        if trace_run_id:
                            _session_run_map[body.sessionId] = trace_run_id

                    # --- Tool progress events → reasoning ---
                    if current_event == "hermes.tool.progress":
                        tool_name = data.get("tool") or "unknown"
                        status = data.get("status") or ""
                        if status:
                            trace_tool_events.append(
                                {
                                    "event": f"tool.{status}",
                                    "tool": tool_name,
                                    "status": status,
                                    "label": data.get("label"),
                                }
                            )
                        if status == "running":
                            label = data.get("label") or tool_name
                            yield _sse_line({
                                **base,
                                "message": "",
                                "reasoningMessage": f"[工具] 调用: {label}",
                                "type": 4,
                            })
                        elif status == "completed":
                            yield _sse_line({
                                **base,
                                "message": "",
                                "reasoningMessage": f"[工具] 完成: {tool_name}",
                                "type": 4,
                            })
                        current_event = ""
                        continue

                    # --- Reasoning/thinking deltas → reasoningMessage ---
                    if current_event == "hermes.reasoning.delta":
                        reasoning_text = data.get("text") or ""
                        if reasoning_text:
                            yield _sse_line({
                                **base,
                                "message": "",
                                "reasoningMessage": reasoning_text,
                                "type": 4,
                            })
                        current_event = ""
                        continue

                    # --- Chat completion chunks → message ---
                    choices = data.get("choices") or []
                    if choices:
                        delta = choices[0].get("delta") or {}
                        content = delta.get("content") or ""
                        if content:
                            full_text += content
                            pending_text += content
                            rewritten_delta = _pop_rewritten_pending()
                            if rewritten_delta:
                                yield _sse_line({**base, "message": rewritten_delta, "reasoningMessage": "", "type": 4})

                        # Check for finish — extract usage from same chunk
                        finish_reason = choices[0].get("finish_reason")
                        if finish_reason:
                            _usage = data.get("usage") or {}
                            _model_name = data.get("model") or ""
                            trace_status = "completed"
                            break

                    current_event = ""

        except httpx.ConnectError as exc:
            logger.error("[%s] hermes stream connect error: %s", log_tag, exc)
            trace_status = "stream_connect_error"
        except Exception as exc:
            logger.error("[%s] hermes stream error: %s", log_tag, exc)
            trace_status = "stream_error"
        finally:
            await event_stream.aclose()
            await client.aclose()

        rewritten_tail = _pop_rewritten_pending(force=True)
        if rewritten_tail:
            yield _sse_line({**base, "message": rewritten_tail, "reasoningMessage": "", "type": 4})

        # Emit token usage event (type=999) before stop
        total_tokens = _usage.get("total_tokens", 0)
        if total_tokens:
            yield _sse_line({**base, "message": {"totalToken": total_tokens}, "reasoningMessage": "", "type": 999})

        if trace_status == "not_started":
            trace_status = "stream_ended"
        _write_trace_once(trace_status)

        # Send stop signal
        yield _sse_line({**base, "message": "[stop]", "reasoningMessage": "", "type": 4})

        # Cleanup: release container, remove run mapping, schedule container mapping cleanup
        _session_run_map.pop(body.sessionId, None)
        await _release_container(base_url, body.sessionId)

        async def _delayed_cleanup(sid: str):
            await asyncio.sleep(_SESSION_CONTAINER_TTL)
            _session_container_map.pop(sid, None)
            logger.debug("[%s] 清理 session→container 映射: %s", log_tag, sid)

        asyncio.create_task(_delayed_cleanup(body.sessionId))
        logger.info("[%s] 请求结束: sessionId=%s", log_tag, body.sessionId)

    return StreamingResponse(
        _stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
            "Connection": "keep-alive",
        },
    )


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

def _build_chat_messages(body: ModelChatRequest) -> list[dict[str, str]]:
    """Convert request messages to the format expected by hermes.

    Handles both plain string content and multi-content arrays with file
    attachments.  Attachment files are downloaded and their text content
    appended to the user message.
    """
    logger.info(f"收到的请求信息:{body}")
    return [{"role": msg.role, "content": _extract_content_text(msg.content)} for msg in body.messages]


def _build_translate_messages(body: ModelChatRequest, target_lang: str) -> list[dict[str, str]]:
    """Build messages with a translation prompt prepended to the last user message."""
    messages = _build_chat_messages(body)
    # Find the last user message and prepend the translation instruction
    for i in range(len(messages) - 1, -1, -1):
        if messages[i]["role"] == "user":
            if target_lang == "zh":
                messages[i]["content"] = f"请将以下内容翻译为中文：\n\n{messages[i]['content']}"
            else:
                messages[i]["content"] = f"Please translate the following content to English:\n\n{messages[i]['content']}"
            break
    return messages


def _build_mindmap_messages(body: ModelChatRequest) -> list[dict[str, str]]:
    """Build messages with a mindmap generation prompt prepended to the last user message."""
    messages = _build_chat_messages(body)
    for i in range(len(messages) - 1, -1, -1):
        if messages[i]["role"] == "user":
            messages[i]["content"] = (
                f"请根据以下内容生成一个markdown格式的思维导图，使用层级标题（# ## ### 等）来组织结构：\n\n{messages[i]['content']}"
            )
            break
    return messages


def _build_sourcepaper_messages(body: ModelChatRequest) -> list[dict[str, str]]:
    """Build messages with a source paper search prompt using medical-keyword-search skill."""
    messages = _build_chat_messages(body)
    for i in range(len(messages) - 1, -1, -1):
        if messages[i]["role"] == "user":
            messages[i]["content"] = (
                f"请根据以下内容，使用 medical-keyword-search skill 检索可能的引用原文文献。"
                f"请以 JSON 格式返回结果，结构如下：\n"
                f'{{"papers": [{{"title": "文献标题", "authors": "作者", "journal": "期刊名", "year": "年份", "abstract": "摘要", "link": "链接"}}]}}\n'
                f"最多返回 5 篇，按相关性排序。\n\n内容：\n\n{messages[i]['content']}"
            )
            break
    return messages


@router.post("/api/science/chat")
async def science_chat(
    request: Request,
    body: ModelChatRequest,
    user=None,
):
    # --- Abort ---
    if body.type == -1:
        await _abort_session(body.sessionId)
        return {"linkId": body.linkId, "sessionId": body.sessionId, "ok": True}

    chat_messages = await asyncio.to_thread(_build_chat_messages, body)
    return await _do_science_chat(request, body, chat_messages)


@router.post("/api/science/translate_to_zh")
async def science_translate_to_zh(
    request: Request,
    body: ModelChatRequest,
    user=None,
):
    """Translate the last user message to Chinese."""
    # --- Abort ---
    if body.type == -1:
        await _abort_session(body.sessionId)
        return {"linkId": body.linkId, "sessionId": body.sessionId, "ok": True}

    chat_messages = await asyncio.to_thread(_build_translate_messages, body, "zh")
    return await _do_science_chat(request, body, chat_messages, log_tag="translate-zh")


@router.post("/api/science/translate_to_en")
async def science_translate_to_en(
    request: Request,
    body: ModelChatRequest,
    user=None,
):
    """Translate the last user message to English."""
    # --- Abort ---
    if body.type == -1:
        await _abort_session(body.sessionId)
        return {"linkId": body.linkId, "sessionId": body.sessionId, "ok": True}

    chat_messages = await asyncio.to_thread(_build_translate_messages, body, "en")
    return await _do_science_chat(request, body, chat_messages, log_tag="translate-en")


@router.post("/api/science/mindmap")
async def science_mindmap(
    request: Request,
    body: ModelChatRequest,
    user=None,
):
    """Generate a markdown mindmap from the last user message."""
    # --- Abort ---
    if body.type == -1:
        await _abort_session(body.sessionId)
        return {"linkId": body.linkId, "sessionId": body.sessionId, "ok": True}

    chat_messages = await asyncio.to_thread(_build_mindmap_messages, body)
    return await _do_science_chat(request, body, chat_messages, log_tag="mindmap")


@router.post("/api/science/sourcepaper")
async def science_sourcepaper(
    request: Request,
    body: ModelChatRequest,
    user=None,
):
    """Search source papers using medical-keyword-search skill, return JSON."""
    # --- Abort ---
    if body.type == -1:
        await _abort_session(body.sessionId)
        return {"linkId": body.linkId, "sessionId": body.sessionId, "ok": True}

    chat_messages = await asyncio.to_thread(_build_sourcepaper_messages, body)
    return await _do_science_chat(request, body, chat_messages, log_tag="sourcepaper")


# ---------------------------------------------------------------------------
# File download endpoint
# ---------------------------------------------------------------------------

@router.get("/api/science/chat/file")
async def science_chat_file(
    sessionId: str = Query(..., description="会话ID"),
    path: str = Query(..., description="容器内文件绝对路径，如 /tmp/report.pdf"),
):
    """从 hermes 容器下载生成的文件（支持任意路径和文件类型）。"""
    container_url = _session_container_map.get(sessionId)
    if not container_url:
        raise HTTPException(status_code=404, detail="会话未找到或已过期")
    container_name = _container_name_from_url(container_url)

    # 安全校验：禁止路径穿越和敏感文件
    normalized = path.replace("\\", "/")
    if ".." in normalized:
        raise HTTPException(status_code=400, detail="非法路径")
    _SENSITIVE = {".env", "config.yaml", "config.yml", ".ssh", "id_rsa", "id_ed25519"}
    if any(seg in _SENSITIVE for seg in normalized.split("/")):
        raise HTTPException(status_code=403, detail="无法访问该文件")

    try:
        container = get_docker_container(container_name)
        stream, _stat = container.get_archive(normalized)
        archive = b"".join(stream)
        with tarfile.open(fileobj=io.BytesIO(archive), mode="r:*") as tar:
            for member in tar:
                if not member.isfile():
                    continue
                extracted = tar.extractfile(member)
                if extracted is None:
                    continue
                content = extracted.read()
                media_type = mimetypes.guess_type(normalized)[0] or "application/octet-stream"
                filename = normalized.rsplit("/", 1)[-1]
                return Response(
                    content=content,
                    media_type=media_type,
                    headers={"Content-Disposition": f'attachment; filename="{filename}"'},
                )
    except DockerNotFound:
        raise HTTPException(status_code=404, detail="文件未找到")
    except DockerAPIError as exc:
        logger.error("[model-chat] file download docker error: %s", exc)
        raise HTTPException(status_code=500, detail="文件读取失败")
    except Exception as exc:
        logger.error("[model-chat] file download error: %s", exc)
        raise HTTPException(status_code=500, detail="文件读取失败")

    raise HTTPException(status_code=404, detail="文件未找到")
